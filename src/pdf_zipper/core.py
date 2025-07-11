"""
Core PDF processing functionality.

This module contains the main PDF compression and conversion logic.
"""

import io
import os
import platform
import tempfile
from pathlib import Path
from typing import Callable, Optional, Tuple

import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

# 支持的文件类型
SUPPORTED_INPUT_TYPES = {
    '.pdf': 'PDF Document',
    '.pptx': 'PowerPoint Presentation',
}

SUPPORTED_OUTPUT_TYPES = {
    '.pdf': 'PDF Document',
    '.pptx': 'PowerPoint Presentation',
}


def get_file_type(file_path: str) -> Tuple[str, str]:
    """
    获取文件类型信息

    Args:
        file_path: 文件路径

    Returns:
        Tuple[extension, description]: 文件扩展名和描述
    """
    path = Path(file_path)
    ext = path.suffix.lower()

    if ext in SUPPORTED_INPUT_TYPES:
        return ext, SUPPORTED_INPUT_TYPES[ext]
    else:
        return ext, "Unknown file type"


def validate_input_file(file_path: str) -> bool:
    """
    验证输入文件是否支持

    Args:
        file_path: 文件路径

    Returns:
        bool: 是否支持该文件类型
    """
    if not os.path.exists(file_path):
        return False

    ext, _ = get_file_type(file_path)
    return ext in SUPPORTED_INPUT_TYPES


def _generate_pdf_data(
    doc: fitz.Document, dpi: float, logger_func: Callable[[str], None]
) -> Optional[bytes]:
    """Generates PDF data in memory and logs progress."""
    image_list = []
    total_pages = len(doc)
    logger_func(f"Converting {total_pages} pages with DPI {int(dpi)}...")

    for i, page in enumerate(doc):
        try:
            pix = page.get_pixmap(dpi=int(dpi))
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            image_list.append(img)
            if (i + 1) % 10 == 0 or i + 1 == total_pages:
                logger_func(f"  - Converted page {i+1}/{total_pages}")
        except Exception as e:
            logger_func(
                f"[bold red]Warning:[/bold red] Error processing page {i+1}, skipped. Error: {e}"
            )
            continue

    if not image_list:
        logger_func(
            "[bold red]Error:[/bold red] No images were generated from the PDF."
        )
        return None

    logger_func("Saving PDF data to memory buffer...")
    pdf_buffer = io.BytesIO()
    image_list[0].save(
        pdf_buffer,
        "PDF",
        resolution=float(dpi),
        save_all=True,
        append_images=image_list[1:],
    )
    logger_func("[green]In-memory save complete.[/green]")
    return pdf_buffer.getvalue()


def _find_optimal_dpi(
    doc: fitz.Document,
    target_size_mb: float,
    logger_func: Callable[[str], None],
    tolerance: float = 0.05,
) -> Tuple[float, Optional[bytes]]:
    """
    使用二分搜索找到最佳DPI以达到目标文件大小

    Args:
        doc: PyMuPDF文档对象
        target_size_mb: 目标文件大小（MB）
        logger_func: 日志记录函数
        tolerance: 容差范围

    Returns:
        Tuple[best_dpi, best_pdf_data]: 最佳DPI和对应的PDF数据
    """
    dpi_low, dpi_high = 30, 300
    best_dpi, best_pdf_data = dpi_low, None

    ITERATIONS = 7
    logger_func(f"Starting iterative search for best DPI ({ITERATIONS} iterations)...")

    for i in range(ITERATIONS):
        dpi_guess = (dpi_low + dpi_high) / 2
        if dpi_guess < dpi_low + 1:
            break

        logger_func(
            f"\n[bold]Attempt {i + 1}/{ITERATIONS}:[/bold] Trying DPI = {int(dpi_guess)}"
        )
        pdf_data = _generate_pdf_data(doc, dpi_guess, logger_func)

        if not pdf_data:
            logger_func(
                f"[yellow]Failed to generate PDF at DPI {int(dpi_guess)}.[/yellow]"
            )
            dpi_high = dpi_guess
            continue

        current_size_mb = len(pdf_data) / (1024 * 1024)
        logger_func(f"  - Generated size: {current_size_mb:.2f} MB")

        if abs(current_size_mb - target_size_mb) / target_size_mb <= tolerance:
            logger_func("[green]Result is within tolerance. Search finished.[/green]")
            best_dpi, best_pdf_data = dpi_guess, pdf_data
            break

        if current_size_mb > target_size_mb:
            dpi_high = dpi_guess
        else:
            dpi_low = dpi_guess
            best_dpi, best_pdf_data = dpi_guess, pdf_data

    return best_dpi, best_pdf_data


def _find_optimal_dpi_for_pptx(
    doc: fitz.Document,
    target_size_mb: float,
    logger_func: Callable[[str], None],
    tolerance: float = 0.05,
) -> Tuple[float, Optional[str]]:
    """
    使用二分搜索找到最佳DPI以达到目标PPTX文件大小

    Args:
        doc: PyMuPDF文档对象
        target_size_mb: 目标PPTX文件大小（MB）
        logger_func: 日志记录函数
        tolerance: 容差范围

    Returns:
        Tuple[best_dpi, best_pptx_path]: 最佳DPI和对应的PPTX文件路径
    """
    dpi_low, dpi_high = 30, 300
    best_dpi, best_pptx_path = dpi_low, None

    ITERATIONS = 7
    logger_func(f"Starting iterative search for best DPI targeting PPTX size ({ITERATIONS} iterations)...")

    for i in range(ITERATIONS):
        dpi_guess = (dpi_low + dpi_high) / 2
        if dpi_guess < dpi_low + 1:
            break

        logger_func(
            f"\n[bold]Attempt {i + 1}/{ITERATIONS}:[/bold] Trying DPI = {int(dpi_guess)}"
        )

        # 生成PDF数据
        pdf_data = _generate_pdf_data(doc, dpi_guess, logger_func)
        if not pdf_data:
            logger_func(
                f"[yellow]Failed to generate PDF at DPI {int(dpi_guess)}.[/yellow]"
            )
            dpi_high = dpi_guess
            continue

        # 将PDF转换为PPTX并检查PPTX文件大小
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as temp_pdf:
            temp_pdf_path = temp_pdf.name
            with open(temp_pdf_path, "wb") as f:
                f.write(pdf_data)

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_pptx:
            temp_pptx_path = temp_pptx.name

        try:
            # 转换为PPTX
            convert_to_ppt(temp_pdf_path, temp_pptx_path, int(dpi_guess), lambda msg: None)  # 静默转换

            if os.path.exists(temp_pptx_path):
                current_size_mb = os.path.getsize(temp_pptx_path) / (1024 * 1024)
                logger_func(f"  - Generated PPTX size: {current_size_mb:.2f} MB")

                if abs(current_size_mb - target_size_mb) / target_size_mb <= tolerance:
                    logger_func("[green]PPTX size is within tolerance. Search finished.[/green]")
                    # 清理之前的最佳文件
                    if best_pptx_path and os.path.exists(best_pptx_path):
                        os.unlink(best_pptx_path)
                    best_dpi, best_pptx_path = dpi_guess, temp_pptx_path
                    break

                if current_size_mb > target_size_mb:
                    dpi_high = dpi_guess
                    # 清理当前文件
                    os.unlink(temp_pptx_path)
                else:
                    dpi_low = dpi_guess
                    # 清理之前的最佳文件
                    if best_pptx_path and os.path.exists(best_pptx_path):
                        os.unlink(best_pptx_path)
                    best_dpi, best_pptx_path = dpi_guess, temp_pptx_path
            else:
                logger_func(f"[yellow]Failed to generate PPTX at DPI {int(dpi_guess)}.[/yellow]")
                dpi_high = dpi_guess

        finally:
            # 清理临时PDF文件
            if os.path.exists(temp_pdf_path):
                os.unlink(temp_pdf_path)

    return best_dpi, best_pptx_path


def autocompress_pdf(
    input_path: str,
    output_path: str,
    target_size_mb: float,
    logger_func: Callable[[str], None],
    tolerance: float = 0.05,
) -> None:
    """Automatically adjusts DPI to compress a PDF to a target size."""
    if not os.path.exists(input_path):
        logger_func(f"[bold red]Error:[/bold red] Input file not found: '{input_path}'")
        return

    original_size_mb = os.path.getsize(input_path) / (1024 * 1024)
    if original_size_mb <= target_size_mb:
        logger_func(
            "[bold yellow]Info:[/bold yellow] Target size is not smaller than the original. Copying file directly."
        )
        import shutil

        shutil.copy(input_path, output_path)
        return

    logger_func("Starting auto-compression...")
    logger_func(f" - Original Size: {original_size_mb:.2f} MB")
    logger_func(f" - Target Size:   {target_size_mb:.2f} MB")

    try:
        doc = fitz.open(input_path)
    except Exception as e:
        logger_func(f"[bold red]Error:[/bold red] Could not open PDF. Reason: {e}")
        return

    best_dpi, best_pdf_data = _find_optimal_dpi(doc, target_size_mb, logger_func, tolerance)
    doc.close()

    if best_pdf_data is None:
        logger_func(
            "[yellow]Search finished without a perfect match. Regenerating with best found DPI...[/yellow]"
        )
        doc = fitz.open(input_path)
        best_pdf_data = _generate_pdf_data(doc, best_dpi, logger_func)
        doc.close()

    if not best_pdf_data:
        logger_func("[bold red]Error:[/bold red] Failed to generate final PDF.")
        return

    logger_func(
        f"\nSearch complete. Best DPI found: [bold cyan]{int(best_dpi)}[/bold cyan]. Saving file..."
    )
    with open(output_path, "wb") as f:
        f.write(best_pdf_data)

    final_size_mb = os.path.getsize(output_path) / (1024 * 1024)
    ratio = (1 - final_size_mb / original_size_mb) * 100 if original_size_mb > 0 else 0

    logger_func("\n[bold green]🎉 Auto-compression Complete![/bold green]")
    logger_func(f"  - Original size: {original_size_mb:.2f} MB")
    logger_func(
        f"  - Final size:    {final_size_mb:.2f} MB (Target: {target_size_mb:.2f} MB)"
    )
    logger_func(f"  - Compression:   {ratio:.2f}%")
    logger_func(f"  - Saved to:      [cyan]{output_path}[/cyan]")


def autocompress_pptx(
    input_path: str,
    output_path: str,
    target_size_mb: float,
    logger_func: Callable[[str], None],
    tolerance: float = 0.05,
) -> None:
    """
    自动调整DPI以将PPTX文件压缩到目标大小

    工作流程：PPTX → PDF → DPI优化 → PPTX
    """
    if not os.path.exists(input_path):
        logger_func(f"[bold red]Error:[/bold red] Input file not found: '{input_path}'")
        return

    original_size_mb = os.path.getsize(input_path) / (1024 * 1024)
    if original_size_mb <= target_size_mb:
        logger_func(
            "[bold yellow]Info:[/bold yellow] Target size is not smaller than the original. Copying file directly."
        )
        import shutil
        shutil.copy(input_path, output_path)
        return

    logger_func("Starting PPTX auto-compression...")
    logger_func(f" - Original Size: {original_size_mb:.2f} MB")
    logger_func(f" - Target Size:   {target_size_mb:.2f} MB")

    # 创建临时文件
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as temp_pdf:
        temp_pdf_path = temp_pdf.name

    try:
        # 步骤1：将PPTX转换为临时PDF
        logger_func("Step 1: Converting PPTX to temporary PDF...")
        convert_pptx_to_pdf(input_path, temp_pdf_path, logger_func)

        if not os.path.exists(temp_pdf_path):
            logger_func("[bold red]Error:[/bold red] Failed to convert PPTX to PDF.")
            return

        # 步骤2：使用DPI优化找到最佳压缩设置
        logger_func("Step 2: Finding optimal DPI for target size...")
        try:
            doc = fitz.open(temp_pdf_path)
        except Exception as e:
            logger_func(f"[bold red]Error:[/bold red] Could not open temporary PDF. Reason: {e}")
            return

        best_dpi, best_pdf_data = _find_optimal_dpi(doc, target_size_mb, logger_func, tolerance)
        doc.close()

        if best_pdf_data is None:
            logger_func(
                "[yellow]Search finished without a perfect match. Regenerating with best found DPI...[/yellow]"
            )
            doc = fitz.open(temp_pdf_path)
            best_pdf_data = _generate_pdf_data(doc, best_dpi, logger_func)
            doc.close()

        if not best_pdf_data:
            logger_func("[bold red]Error:[/bold red] Failed to generate optimized PDF.")
            return

        # 步骤3：将优化后的PDF转换回PPTX
        logger_func("Step 3: Converting optimized PDF back to PPTX...")
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as optimized_pdf:
            optimized_pdf_path = optimized_pdf.name
            with open(optimized_pdf_path, "wb") as f:
                f.write(best_pdf_data)

        try:
            convert_to_ppt(optimized_pdf_path, output_path, int(best_dpi), logger_func)

            if os.path.exists(output_path):
                final_size_mb = os.path.getsize(output_path) / (1024 * 1024)
                ratio = (1 - final_size_mb / original_size_mb) * 100 if original_size_mb > 0 else 0

                logger_func("\n[bold green]🎉 PPTX Auto-compression Complete![/bold green]")
                logger_func(f"  - Original size: {original_size_mb:.2f} MB")
                logger_func(f"  - Final size:    {final_size_mb:.2f} MB (Target: {target_size_mb:.2f} MB)")
                logger_func(f"  - Compression:   {ratio:.2f}%")
                logger_func(f"  - Best DPI:      {int(best_dpi)}")
                logger_func(f"  - Saved to:      [cyan]{output_path}[/cyan]")
            else:
                logger_func("[bold red]Error:[/bold red] Failed to create final PPTX file.")

        finally:
            # 清理临时的优化PDF文件
            if os.path.exists(optimized_pdf_path):
                os.unlink(optimized_pdf_path)

    finally:
        # 清理临时PDF文件
        if os.path.exists(temp_pdf_path):
            os.unlink(temp_pdf_path)


def autocompress_pdf_to_pptx(
    input_path: str,
    output_path: str,
    target_size_mb: float,
    logger_func: Callable[[str], None],
    tolerance: float = 0.05,
) -> None:
    """
    自动调整DPI以将PDF文件压缩到目标大小并输出为PPTX格式

    工作流程：PDF → DPI优化（基于PPTX大小） → PPTX
    """
    if not os.path.exists(input_path):
        logger_func(f"[bold red]Error:[/bold red] Input file not found: '{input_path}'")
        return

    # 验证输入文件是PDF
    ext, _ = get_file_type(input_path)
    if ext != '.pdf':
        logger_func(f"[bold red]Error:[/bold red] Input file must be PDF, got: {ext}")
        return

    original_size_mb = os.path.getsize(input_path) / (1024 * 1024)
    logger_func("Starting PDF to target-size PPTX conversion...")
    logger_func(f" - Original Size: {original_size_mb:.2f} MB")
    logger_func(f" - Target Size:   {target_size_mb:.2f} MB")

    try:
        doc = fitz.open(input_path)
    except Exception as e:
        logger_func(f"[bold red]Error:[/bold red] Could not open PDF. Reason: {e}")
        return

    # 使用新的DPI优化逻辑，直接基于PPTX文件大小进行优化
    logger_func("Finding optimal DPI based on final PPTX size...")
    best_dpi, best_pptx_path = _find_optimal_dpi_for_pptx(doc, target_size_mb, logger_func, tolerance)
    doc.close()

    if best_pptx_path is None or not os.path.exists(best_pptx_path):
        logger_func("[bold red]Error:[/bold red] Failed to generate target-size PPTX.")
        return

    try:
        # 移动最佳PPTX文件到目标位置
        import shutil
        shutil.move(best_pptx_path, output_path)

        final_size_mb = os.path.getsize(output_path) / (1024 * 1024)

        logger_func("\n[bold green]🎉 PDF to Target-size PPTX Complete![/bold green]")
        logger_func(f"  - Original size: {original_size_mb:.2f} MB")
        logger_func(f"  - Final size:    {final_size_mb:.2f} MB (Target: {target_size_mb:.2f} MB)")
        logger_func(f"  - Best DPI:      {int(best_dpi)}")
        logger_func(f"  - Saved to:      [cyan]{output_path}[/cyan]")

    except Exception as e:
        logger_func(f"[bold red]Error:[/bold red] Failed to save final PPTX file: {e}")
        # 清理临时文件
        if best_pptx_path and os.path.exists(best_pptx_path):
            os.unlink(best_pptx_path)


def autocompress(
    input_path: str,
    output_path: str,
    target_size_mb: float,
    logger_func: Callable[[str], None],
    tolerance: float = 0.05,
) -> None:
    """
    通用的自动压缩函数，根据文件类型和输出格式调用相应的压缩函数

    支持的转换：
    - PDF → PDF: 直接压缩
    - PPTX → PPTX: 转换为PDF → 压缩 → 转换回PPTX
    - PDF → PPTX: DPI优化 → 转换为PPTX
    """
    if not validate_input_file(input_path):
        logger_func(f"[bold red]Error:[/bold red] Unsupported file type or file not found: '{input_path}'")
        return

    input_ext, _ = get_file_type(input_path)
    output_ext = Path(output_path).suffix.lower()

    # 根据输入和输出格式选择处理方式
    if input_ext == '.pdf' and output_ext == '.pdf':
        autocompress_pdf(input_path, output_path, target_size_mb, logger_func, tolerance)
    elif input_ext == '.pptx' and output_ext == '.pptx':
        autocompress_pptx(input_path, output_path, target_size_mb, logger_func, tolerance)
    elif input_ext == '.pdf' and output_ext == '.pptx':
        autocompress_pdf_to_pptx(input_path, output_path, target_size_mb, logger_func, tolerance)
    else:
        logger_func(f"[bold red]Error:[/bold red] Unsupported conversion: {input_ext} → {output_ext}")
        logger_func("Supported conversions: PDF→PDF, PPTX→PPTX, PDF→PPTX")


def compress_pdf(
    input_path: str, output_path: str, dpi: int, logger_func: Callable[[str], None]
) -> None:
    """Manually compresses a PDF using a specific DPI."""
    logger_func("Starting manual compression...")
    doc = fitz.open(input_path)
    pdf_data = _generate_pdf_data(doc, dpi, logger_func)
    doc.close()

    if pdf_data:
        with open(output_path, "wb") as f:
            f.write(pdf_data)

        original_size = os.path.getsize(input_path) / (1024 * 1024)
        final_size = os.path.getsize(output_path) / (1024 * 1024)
        ratio = (1 - final_size / original_size) * 100 if original_size > 0 else 0

        logger_func("\n[bold green]🎉 Manual Compression Complete![/bold green]")
        logger_func(f"  - Original size: {original_size:.2f} MB")
        logger_func(f"  - Final size:    {final_size:.2f} MB")
        logger_func(f"  - Compression:   {ratio:.2f}%")
        logger_func(f"  - Saved to:      [cyan]{output_path}[/cyan]")


def check_conversion_tools() -> dict:
    """检查系统中可用的转换工具"""
    tools = {}

    try:
        import subprocess

        # 检查 LibreOffice
        try:
            subprocess.run(["libreoffice", "--version"],
                         capture_output=True, timeout=5, check=True)
            tools["libreoffice"] = True
        except (FileNotFoundError, subprocess.CalledProcessError, subprocess.TimeoutExpired):
            tools["libreoffice"] = False

        # 检查 unoconv
        try:
            subprocess.run(["unoconv", "--version"],
                         capture_output=True, timeout=5, check=True)
            tools["unoconv"] = True
        except (FileNotFoundError, subprocess.CalledProcessError, subprocess.TimeoutExpired):
            tools["unoconv"] = False



    except ImportError:
        # subprocess 不可用
        tools = {"libreoffice": False, "unoconv": False}

    return tools


def convert_pptx_to_pdf(
    input_path: str, output_path: str, logger_func: Callable[[str], None]
) -> None:
    """
    将 PowerPoint 文件转换为 PDF

    Args:
        input_path: 输入的 PPTX 文件路径
        output_path: 输出的 PDF 文件路径
        logger_func: 日志记录函数
    """
    logger_func("Starting PPTX to PDF conversion...")

    if not os.path.exists(input_path):
        logger_func(f"[bold red]Error:[/bold red] Input file not found: '{input_path}'")
        return

    # 检查可用的转换工具
    available_tools = check_conversion_tools()
    has_good_tools = any(available_tools.values())

    if not has_good_tools:
        logger_func("[bold yellow]Warning:[/bold yellow] No advanced conversion tools found.")
        logger_func("For better PPTX to PDF conversion with full styling, consider installing:")
        logger_func("  - LibreOffice: https://www.libreoffice.org/")
        logger_func("  - Or unoconv: pip install unoconv")
        logger_func("Using fallback method (text-only)...")

    try:
        # 使用不同平台的转换方法
        if platform.system() == "Windows":
            _convert_pptx_to_pdf_windows(input_path, output_path, logger_func)
        else:
            _convert_pptx_to_pdf_cross_platform(input_path, output_path, logger_func)

        # 检查输出文件是否成功创建
        if os.path.exists(output_path):
            input_size = os.path.getsize(input_path) / (1024 * 1024)
            output_size = os.path.getsize(output_path) / (1024 * 1024)

            logger_func("\n[bold green]🎉 PPTX to PDF Conversion Complete![/bold green]")
            logger_func(f"  - Input size:  {input_size:.2f} MB")
            logger_func(f"  - Output size: {output_size:.2f} MB")
            logger_func(f"  - Saved to:    [cyan]{output_path}[/cyan]")
        else:
            logger_func("[bold red]Error:[/bold red] Failed to create output PDF file.")

    except Exception as e:
        logger_func(f"[bold red]Error:[/bold red] PPTX to PDF conversion failed: {e}")


def _convert_pptx_to_pdf_windows(
    input_path: str, output_path: str, logger_func: Callable[[str], None]
) -> None:
    """Windows 平台使用 COM 接口转换 PPTX 到 PDF"""
    try:
        import comtypes.client

        logger_func("Using Windows PowerPoint COM interface...")

        # 创建 PowerPoint 应用实例
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1

        # 打开演示文稿
        presentation = powerpoint.Presentations.Open(os.path.abspath(input_path))

        # 导出为 PDF (格式代码 32 表示 PDF)
        presentation.SaveAs(os.path.abspath(output_path), 32)

        # 关闭演示文稿和应用
        presentation.Close()
        powerpoint.Quit()

        logger_func("✅ Windows COM conversion completed")

    except ImportError:
        logger_func("[bold yellow]Warning:[/bold yellow] Windows COM libraries not available, falling back to cross-platform method")
        _convert_pptx_to_pdf_cross_platform(input_path, output_path, logger_func)
    except Exception as e:
        logger_func(f"[bold red]Error:[/bold red] Windows COM conversion failed: {e}")
        logger_func("Falling back to cross-platform method...")
        _convert_pptx_to_pdf_cross_platform(input_path, output_path, logger_func)


def _convert_pptx_to_pdf_cross_platform(
    input_path: str, output_path: str, logger_func: Callable[[str], None]
) -> None:
    """跨平台方法：优先使用系统工具，保持完整样式"""
    logger_func("Using cross-platform conversion method...")

    # 尝试多种系统工具进行转换
    conversion_methods = [
        ("LibreOffice", _try_libreoffice_conversion),
        ("unoconv", _try_unoconv_conversion),
    ]

    for method_name, conversion_func in conversion_methods:
        logger_func(f"Trying {method_name} conversion...")
        try:
            if conversion_func(input_path, output_path, logger_func):
                logger_func(f"✅ {method_name} conversion successful!")
                return
        except Exception as e:
            logger_func(f"⚠️  {method_name} conversion failed: {e}")
            continue

    # 如果所有系统工具都失败，使用备用方法
    logger_func("All system tools failed, using fallback method...")
    _convert_pptx_fallback_method(input_path, output_path, logger_func)


def _try_libreoffice_conversion(input_pptx: str, output_pdf: str, logger_func: Callable[[str], None]) -> bool:
    """使用 LibreOffice 转换 PPTX 到 PDF（保持完整样式）"""
    try:
        import subprocess

        # 确保输出目录存在
        output_dir = os.path.dirname(output_pdf)
        os.makedirs(output_dir, exist_ok=True)

        # 使用 LibreOffice 命令行转换
        cmd = [
            "libreoffice",
            "--headless",
            "--convert-to", "pdf",
            "--outdir", output_dir,
            input_pptx
        ]

        logger_func("  - Running LibreOffice conversion...")
        subprocess.run(cmd, check=True, capture_output=True, timeout=60, text=True)

        # LibreOffice 会创建与输入文件同名但扩展名为 .pdf 的文件
        expected_output = os.path.join(
            output_dir,
            os.path.splitext(os.path.basename(input_pptx))[0] + ".pdf"
        )

        if os.path.exists(expected_output):
            if expected_output != output_pdf:
                os.rename(expected_output, output_pdf)
            return True

    except subprocess.TimeoutExpired:
        logger_func("  - LibreOffice conversion timed out")
    except subprocess.CalledProcessError as e:
        logger_func(f"  - LibreOffice conversion failed: {e.stderr}")
    except FileNotFoundError:
        logger_func("  - LibreOffice not found in system PATH")
    except Exception as e:
        logger_func(f"  - LibreOffice conversion error: {e}")

    return False


def _try_unoconv_conversion(input_pptx: str, output_pdf: str, logger_func: Callable[[str], None]) -> bool:
    """使用 unoconv 转换 PPTX 到 PDF"""
    try:
        import subprocess

        logger_func("  - Running unoconv conversion...")
        cmd = ["unoconv", "-f", "pdf", "-o", output_pdf, input_pptx]

        subprocess.run(cmd, check=True, capture_output=True, timeout=60, text=True)
        return os.path.exists(output_pdf)

    except subprocess.TimeoutExpired:
        logger_func("  - unoconv conversion timed out")
    except subprocess.CalledProcessError as e:
        logger_func(f"  - unoconv conversion failed: {e.stderr}")
    except FileNotFoundError:
        logger_func("  - unoconv not found in system PATH")
    except Exception as e:
        logger_func(f"  - unoconv conversion error: {e}")

    return False



def _convert_pptx_fallback_method(input_path: str, output_path: str, logger_func: Callable[[str], None]) -> None:
    """备用方法：创建包含幻灯片信息的 PDF"""
    try:
        logger_func("Using fallback method (text-only conversion)...")

        # 读取 PowerPoint 文件
        prs = Presentation(input_path)
        total_slides = len(prs.slides)
        logger_func(f"Found {total_slides} slides to convert")

        if total_slides == 0:
            logger_func("[bold red]Error:[/bold red] No slides found in presentation")
            return

        # 创建 PDF 文档
        doc = fitz.open()

        for i, slide in enumerate(prs.slides):
            logger_func(f"  - Processing slide {i+1}/{total_slides}")

            # 创建新页面
            page = doc.new_page()

            # 添加标题
            title_rect = fitz.Rect(50, 50, 550, 100)
            page.insert_textbox(title_rect, f"Slide {i+1}", fontsize=20, align=1)

            # 提取并添加文本内容
            y_position = 120
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text_rect = fitz.Rect(50, y_position, 550, y_position + 50)
                    page.insert_textbox(text_rect, shape.text.strip(), fontsize=12, align=0)
                    y_position += 60

                    if y_position > 750:  # 避免超出页面
                        break

            # 添加警告信息
            warning_rect = fitz.Rect(50, 750, 550, 780)
            page.insert_textbox(
                warning_rect,
                "⚠️ This is a text-only conversion. For full styling, install LibreOffice.",
                fontsize=10,
                align=1
            )

        # 保存 PDF
        doc.save(output_path)
        doc.close()

        logger_func(f"✅ Fallback conversion completed with {total_slides} pages")
        logger_func("[bold yellow]Note:[/bold yellow] This conversion preserves text only.")
        logger_func("[bold yellow]Tip:[/bold yellow] Install LibreOffice for full styling support:")
        logger_func("  - macOS: brew install --cask libreoffice")
        logger_func("  - Ubuntu: sudo apt install libreoffice")
        logger_func("  - Windows: Download from https://www.libreoffice.org/")

    except Exception as e:
        logger_func(f"[bold red]Error:[/bold red] Fallback conversion failed: {e}")


def _create_placeholder_pdf(output_path: str, text: str, logger_func: Callable[[str], None]) -> None:
    """创建一个包含文本的占位符 PDF"""
    try:
        doc = fitz.open()  # 创建新的 PDF 文档
        page = doc.new_page()  # 添加新页面

        # 在页面中心添加文本
        rect = page.rect
        text_rect = fitz.Rect(50, rect.height/2 - 50, rect.width - 50, rect.height/2 + 50)
        page.insert_textbox(text_rect, text, fontsize=24, align=1)  # 居中对齐

        doc.save(output_path)
        doc.close()

    except Exception as e:
        logger_func(f"[bold yellow]Warning:[/bold yellow] Failed to create placeholder PDF: {e}")


def _merge_pdfs(pdf_paths: list, output_path: str, logger_func: Callable[[str], None]) -> None:
    """合并多个 PDF 文件"""
    try:
        merged_doc = fitz.open()

        for pdf_path in pdf_paths:
            if os.path.exists(pdf_path):
                doc = fitz.open(pdf_path)
                merged_doc.insert_pdf(doc)
                doc.close()

        merged_doc.save(output_path)
        merged_doc.close()

        logger_func(f"✅ Merged {len(pdf_paths)} pages into final PDF")

    except Exception as e:
        logger_func(f"[bold red]Error:[/bold red] Failed to merge PDFs: {e}")


def _optimize_image_for_pptx(img: Image.Image, target_quality: int = 85) -> bytes:
    """
    优化图像以减小PPTX文件大小

    Args:
        img: PIL图像对象
        target_quality: JPEG质量 (1-100)

    Returns:
        bytes: 压缩后的图像数据
    """
    img_buffer = io.BytesIO()

    # 如果图像很大，先进行适度缩放
    max_dimension = 1920  # 最大尺寸限制
    if max(img.width, img.height) > max_dimension:
        ratio = max_dimension / max(img.width, img.height)
        new_width = int(img.width * ratio)
        new_height = int(img.height * ratio)
        img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)

    # 使用JPEG压缩，在质量和文件大小之间取得平衡
    img.save(img_buffer, format='JPEG', quality=target_quality, optimize=True)
    return img_buffer.getvalue()


def convert_to_ppt(
    input_path: str, output_path: str, dpi: int, logger_func: Callable[[str], None]
) -> None:
    """Converts a PDF to a PowerPoint presentation with optimized image compression."""
    logger_func("Starting PDF to PPT conversion...")
    doc = fitz.open(input_path)
    prs = Presentation()
    total_pages = len(doc)

    # 根据页面数量调整压缩质量
    if total_pages > 50:
        quality = 70  # 页面很多时使用更高压缩
        logger_func(f"Large document ({total_pages} pages), using higher compression (quality: {quality})")
    elif total_pages > 20:
        quality = 80  # 中等页面数使用中等压缩
        logger_func(f"Medium document ({total_pages} pages), using medium compression (quality: {quality})")
    else:
        quality = 85  # 页面较少时使用较低压缩
        logger_func(f"Small document ({total_pages} pages), using standard compression (quality: {quality})")

    for i, page in enumerate(doc):
        logger_func(f"  - Processing page {i+1}/{total_pages}")
        pix = page.get_pixmap(dpi=int(dpi))

        # 转换为PIL图像并优化
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        img_data = _optimize_image_for_pptx(img, quality)

        prs.slide_width = Inches(page.rect.width / 72)
        prs.slide_height = Inches(page.rect.height / 72)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(
            io.BytesIO(img_data), 0, 0, width=prs.slide_width, height=prs.slide_height
        )

    doc.close()
    prs.save(output_path)
    logger_func("\n[bold green]🎉 Conversion Complete![/bold green]")
    logger_func(f"  - Saved to: [cyan]{output_path}[/cyan]")
