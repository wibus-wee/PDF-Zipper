import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
import os
import io
from pathlib import Path

from textual.app import App, ComposeResult
from textual.containers import Container, Vertical, Horizontal
from textual.widgets import (
    Button,
    DirectoryTree,
    Footer,
    Header,
    Input,
    Label,
    Static,
    TabbedContent,
    TabPane,
    RichLog,
)
from textual.validation import Number
from textual.events import Paste
from textual import work


# --- PDF Core Logic ---
# Note: These functions are slightly modified to log to the Textual UI
# instead of printing to the console.

def _generate_pdf_data(doc, dpi, logger_func):
    """Generates PDF data in memory and logs progress."""
    image_list = []
    total_pages = len(doc)
    logger_func(f"Converting {total_pages} pages with DPI {int(dpi)}...")
    for i, page in enumerate(doc):
        try:
            pix = page.get_pixmap(dpi=int(dpi))
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            image_list.append(img)
            if (i + 1) % 10 == 0 or i + 1 == total_pages:
                logger_func(f"  - Converted page {i+1}/{total_pages}")
        except Exception as e:
            logger_func(f"[bold red]Warning:[/bold red] Error processing page {i+1}, skipped. Error: {e}")
            continue
    
    if not image_list:
        logger_func("[bold red]Error:[/bold red] No images were generated from the PDF.")
        return None

    logger_func("Saving PDF data to memory buffer...")
    pdf_buffer = io.BytesIO()
    image_list[0].save(
        pdf_buffer, "PDF", resolution=float(dpi), save_all=True, append_images=image_list[1:]
    )
    logger_func("[green]In-memory save complete.[/green]")
    return pdf_buffer.getvalue()


def autocompress_pdf(input_path, output_path, target_size_mb, logger_func, tolerance=0.05):
    """Automatically adjusts DPI to compress a PDF to a target size."""
    if not os.path.exists(input_path):
        logger_func(f"[bold red]Error:[/bold red] Input file not found: '{input_path}'")
        return

    original_size_mb = os.path.getsize(input_path) / (1024 * 1024)
    if original_size_mb <= target_size_mb:
        logger_func("[bold yellow]Info:[/bold yellow] Target size is not smaller than the original. Copying file directly.")
        import shutil
        shutil.copy(input_path, output_path)
        return

    logger_func(f"Starting auto-compression...")
    logger_func(f" - Original Size: {original_size_mb:.2f} MB")
    logger_func(f" - Target Size:   {target_size_mb:.2f} MB")

    try:
        doc = fitz.open(input_path)
    except Exception as e:
        logger_func(f"[bold red]Error:[/bold red] Could not open PDF. Reason: {e}")
        return

    dpi_low, dpi_high = 30, 300
    best_dpi, best_pdf_data = dpi_low, None
    
    ITERATIONS = 7
    logger_func(f"Starting iterative search for best DPI ({ITERATIONS} iterations)...")

    for i in range(ITERATIONS):
        dpi_guess = (dpi_low + dpi_high) / 2
        if dpi_guess < dpi_low + 1: break

        logger_func(f"\n[bold]Attempt {i + 1}/{ITERATIONS}:[/bold] Trying DPI = {int(dpi_guess)}")
        pdf_data = _generate_pdf_data(doc, dpi_guess, logger_func)
        
        if not pdf_data:
            logger_func(f"[yellow]Failed to generate PDF at DPI {int(dpi_guess)}.[/yellow]")
            dpi_high = dpi_guess
            continue

        current_size_mb = len(pdf_data) / (1024 * 1024)
        logger_func(f"  - Generated size: {current_size_mb:.2f} MB")

        if abs(current_size_mb - target_size_mb) / target_size_mb <= tolerance:
            logger_func("[green]Result is within tolerance. Search finished.[/green]")
            best_dpi, best_pdf_data = dpi_guess, pdf_data
            break

        if current_size_mb > target_size_mb:
            dpi_high = dpi_guess
        else:
            dpi_low = dpi_guess
            best_dpi, best_pdf_data = dpi_guess, pdf_data
    
    doc.close()

    if best_pdf_data is None:
        logger_func("[yellow]Search finished without a perfect match. Regenerating with best found DPI...[/yellow]")
        doc = fitz.open(input_path)
        best_pdf_data = _generate_pdf_data(doc, best_dpi, logger_func)
        doc.close()

    if not best_pdf_data:
        logger_func("[bold red]Error:[/bold red] Failed to generate final PDF.")
        return

    logger_func(f"\nSearch complete. Best DPI found: [bold cyan]{int(best_dpi)}[/bold cyan]. Saving file...")
    with open(output_path, "wb") as f:
        f.write(best_pdf_data)

    final_size_mb = os.path.getsize(output_path) / (1024 * 1024)
    ratio = (1 - final_size_mb / original_size_mb) * 100 if original_size_mb > 0 else 0
    
    logger_func("\n[bold green]🎉 Auto-compression Complete![/bold green]")
    logger_func(f"  - Original size: {original_size_mb:.2f} MB")
    logger_func(f"  - Final size:    {final_size_mb:.2f} MB (Target: {target_size_mb:.2f} MB)")
    logger_func(f"  - Compression:   {ratio:.2f}%")
    logger_func(f"  - Saved to:      [cyan]{output_path}[/cyan]")


def compress_pdf(input_path, output_path, dpi, logger_func):
    """Manually compresses a PDF using a specific DPI."""
    logger_func("Starting manual compression...")
    doc = fitz.open(input_path)
    pdf_data = _generate_pdf_data(doc, dpi, logger_func)
    doc.close()
    if pdf_data:
        with open(output_path, "wb") as f:
            f.write(pdf_data)
        
        original_size = os.path.getsize(input_path) / (1024 * 1024)
        final_size = os.path.getsize(output_path) / (1024 * 1024)
        ratio = (1 - final_size / original_size) * 100 if original_size > 0 else 0
        logger_func("\n[bold green]🎉 Manual Compression Complete![/bold green]")
        logger_func(f"  - Original size: {original_size:.2f} MB")
        logger_func(f"  - Final size:    {final_size:.2f} MB")
        logger_func(f"  - Compression:   {ratio:.2f}%")
        logger_func(f"  - Saved to:      [cyan]{output_path}[/cyan]")


def convert_to_ppt(input_path, output_path, dpi, logger_func):
    """Converts a PDF to a PowerPoint presentation."""
    logger_func("Starting PDF to PPT conversion...")
    doc = fitz.open(input_path)
    prs = Presentation()
    total_pages = len(doc)
    for i, page in enumerate(doc):
        logger_func(f"  - Processing page {i+1}/{total_pages}")
        pix = page.get_pixmap(dpi=int(dpi))
        img_data = pix.tobytes("png")
        
        prs.slide_width = Inches(page.rect.width / 72)
        prs.slide_height = Inches(page.rect.height / 72)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(io.BytesIO(img_data), 0, 0, width=prs.slide_width, height=prs.slide_height)
            
    doc.close()
    prs.save(output_path)
    logger_func("\n[bold green]🎉 Conversion Complete![/bold green]")
    logger_func(f"  - Saved to: [cyan]{output_path}[/cyan]")


# --- Textual UI ---

class PDFZipperApp(App):
    CSS_PATH = "compress.css"
    BINDINGS = [("q", "quit", "Quit")]

    def compose(self) -> ComposeResult:
        yield Header(name="🚀 PDF Zipper")
        with Container():
            yield DirectoryTree(".", id="tree-view")
            with Vertical(id="main-view"):
                yield Static("Select a PDF file from the tree or enter custom path.", id="selected-file")
                with TabbedContent(id="tabs"):
                    with TabPane("Custom Path", id="tab-custom"):
                        yield Label("PDF File Path:")
                        yield Input(placeholder="输入完整PDF文件路径，如: /Users/name/file.pdf", id="custom-path")
                        with Horizontal(id="custom-operations"):
                            yield Button("Auto (5MB)", variant="primary", id="btn-custom-auto")
                            yield Button("Manual (150DPI)", variant="primary", id="btn-custom-manual")
                            yield Button("To PPT (150DPI)", variant="primary", id="btn-custom-ppt")
                    with TabPane("Auto Compress", id="tab-auto"):
                        yield Label("Target Size (MB):")
                        yield Input(placeholder="e.g., 5.5", id="auto-size", validators=[Number(minimum=0)])
                        yield Button("Start Auto Compress", variant="primary", id="btn-auto")
                    with TabPane("Manual Compress", id="tab-manual"):
                        yield Label("Resolution (DPI):")
                        yield Input(value="150", id="manual-dpi", validators=[Number(minimum=10)])
                        yield Button("Start Manual Compress", variant="primary", id="btn-manual")
                    with TabPane("Convert to PPT", id="tab-ppt"):
                        yield Label("Image Resolution (DPI):")
                        yield Input(value="150", id="ppt-dpi", validators=[Number(minimum=10)])
                        yield Button("Start Conversion", variant="primary", id="btn-ppt")
                yield RichLog(id="log", wrap=True, highlight=True)
        yield Footer()

    def on_mount(self) -> None:
        """Focus the directory tree on startup."""
        self.query_one(DirectoryTree).focus()
        self.query_one(RichLog).write("Welcome to PDF Zipper!")
        self.query_one(RichLog).write("💡 提示：")
        self.query_one(RichLog).write("  - 从左侧文件树选择PDF文件")
        self.query_one(RichLog).write("  - 或使用「Custom Path」选项卡输入/拖拽文件路径")
        self.query_one(RichLog).write("  - 然后点击对应的操作按钮开始处理")

    def on_directory_tree_file_selected(self, event: DirectoryTree.FileSelected) -> None:
        """Handle file selection."""
        if str(event.path).lower().endswith(".pdf"):
            self.query_one("#selected-file").update(f"Selected: [bold cyan]{event.path}[/bold cyan]")
            self.selected_path = event.path
            # 同时更新自定义路径输入框
            self.query_one("#custom-path").value = str(event.path)
        else:
            self.query_one("#selected-file").update("[bold red]Please select a .pdf file.[/bold red]")
            self.selected_path = None

    def on_paste(self, event: Paste) -> None:
        """Handle paste events for drag and drop functionality."""
        # 检查粘贴的内容是否是文件路径
        pasted_text = event.text.strip()
        if pasted_text and (pasted_text.lower().endswith('.pdf') or os.path.exists(pasted_text)):
            # 如果当前焦点在自定义路径输入框
            try:
                custom_input = self.query_one("#custom-path")
                if custom_input.has_focus:
                    custom_input.value = pasted_text
                    self._update_selected_path_from_custom()
            except:
                pass

    def _update_selected_path_from_custom(self):
        """Update selected path from custom input."""
        custom_path = self.query_one("#custom-path").value.strip()
        if custom_path and os.path.exists(custom_path) and custom_path.lower().endswith('.pdf'):
            self.selected_path = Path(custom_path)
            self.query_one("#selected-file").update(f"Custom: [bold cyan]{custom_path}[/bold cyan]")
        elif custom_path:
            self.query_one("#selected-file").update(f"[bold red]Invalid path: {custom_path}[/bold red]")
            self.selected_path = None
        else:
            self.query_one("#selected-file").update("Select a PDF file from the tree or enter custom path.")
            self.selected_path = None

    def on_input_changed(self, event: Input.Changed) -> None:
        """Handle input changes."""
        if event.input.id == "custom-path":
            self._update_selected_path_from_custom()

    @work(thread=True)
    def worker_autocompress(self, input_path: str, output_path: str, target_size: float, logger_func) -> None:
        """Worker for auto compression."""
        autocompress_pdf(input_path, output_path, target_size, logger_func)

    @work(thread=True)
    def worker_manual_compress(self, input_path: str, output_path: str, dpi: int, logger_func) -> None:
        """Worker for manual compression."""
        compress_pdf(input_path, output_path, dpi, logger_func)

    @work(thread=True)
    def worker_convert_ppt(self, input_path: str, output_path: str, dpi: int, logger_func) -> None:
        """Worker for PPT conversion."""
        convert_to_ppt(input_path, output_path, dpi, logger_func)

    def on_button_pressed(self, event: Button.Pressed) -> None:
        """Handle button presses and dispatch jobs."""
        log = self.query_one(RichLog)

        # Handle browse button
        if event.button.id == "btn-browse":
            log.write("💡 提示：请在「PDF File Path」输入框中输入完整的PDF文件路径")
            log.write("   例如：/Users/username/Documents/file.pdf")
            log.write("   或者直接拖拽PDF文件到输入框中")
            return

        log.clear()

        if not hasattr(self, "selected_path") or self.selected_path is None:
            log.write("[bold red]Error: No PDF file selected.[/bold red]")
            log.write("请先选择PDF文件：")
            log.write("  1. 从左侧文件树选择PDF文件")
            log.write("  2. 或在「Custom Path」选项卡中输入文件路径")
            return

        input_path = str(self.selected_path)
        base, _ = os.path.splitext(os.path.basename(input_path))

        # Common logger function
        def logger(message):
            self.call_from_thread(log.write, message)

        # Handle custom path operations (quick with defaults)
        if event.button.id == "btn-custom-auto":
            target_size = 5.0  # 默认5MB
            log.write(f"使用默认目标大小: {target_size} MB")
            output_path = f"{base}_auto_compressed.pdf"
            self.worker_autocompress(input_path, output_path, target_size, logger)

        elif event.button.id == "btn-custom-manual":
            dpi = 150  # 默认DPI
            log.write(f"使用默认DPI: {dpi}")
            output_path = f"{base}_manual_compressed.pdf"
            self.worker_manual_compress(input_path, output_path, dpi, logger)

        elif event.button.id == "btn-custom-ppt":
            dpi = 150  # 默认DPI
            log.write(f"使用默认DPI: {dpi}")
            output_path = f"{base}_converted.pptx"
            self.worker_convert_ppt(input_path, output_path, dpi, logger)

        elif event.button.id == "btn-auto":
            target_size_input = self.query_one("#auto-size", Input)
            if not target_size_input.is_valid or not target_size_input.value:
                log.write("[bold red]Error: Invalid target size.[/bold red]")
                return
            target_size = float(target_size_input.value)
            output_path = f"{base}_auto_compressed.pdf"
            self.worker_autocompress(input_path, output_path, target_size, logger)

        elif event.button.id == "btn-manual":
            dpi_input = self.query_one("#manual-dpi", Input)
            if not dpi_input.is_valid or not dpi_input.value:
                log.write("[bold red]Error: Invalid DPI value.[/bold red]")
                return
            dpi = int(dpi_input.value)
            output_path = f"{base}_manual_compressed.pdf"
            self.worker_manual_compress(input_path, output_path, dpi, logger)

        elif event.button.id == "btn-ppt":
            dpi_input = self.query_one("#ppt-dpi", Input)
            if not dpi_input.is_valid or not dpi_input.value:
                log.write("[bold red]Error: Invalid DPI value.[/bold red]")
                return
            dpi = int(dpi_input.value)
            output_path = f"{base}_converted.pptx"
            self.worker_convert_ppt(input_path, output_path, dpi, logger)

if __name__ == "__main__":
    app = PDFZipperApp()
    app.run()